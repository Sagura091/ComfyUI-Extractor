"""
Enhanced PowerPoint Processing with Revolutionary OCR
Extracts images with full slide context, OCR text recognition,
and comprehensive slide content analysis.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from PIL import Image, ImageDraw
import io
import json
import logging
from typing import Dict, List, Any, Optional
import numpy as np

from .revolutionary_ocr import create_revolutionary_ocr
from .caption_api import auto_caption
from .settings import OUT_DIR
from .ocr_config import get_revolutionary_ocr_config

# Setup logging
logger = logging.getLogger(__name__)

class EnhancedPPTXProcessor:
    """Enhanced PowerPoint processor with Revolutionary OCR capabilities"""
    
    def __init__(self, ocr_config: Dict[str, Any] = None):
        # Initialize revolutionary OCR
        self.ocr = create_revolutionary_ocr(ocr_config or {
            'languages': ['en'],
            'enable_table_recognition': True,
            'enable_formula_recognition': True,
            'enable_handwriting': True,
            'confidence_threshold': 0.7
        })
        
        self.logger = logging.getLogger(__name__)
    
    def process_pptx_advanced(self, path: Path, push_record_func) -> Dict[str, Any]:
        """
        Process PowerPoint with advanced OCR and slide analysis
        """
        try:
            prs = Presentation(path)
            doc_id = path.stem
            
            processing_results = {
                'document_id': doc_id,
                'total_slides': len(prs.slides),
                'processed_slides': 0,
                'images': [],
                'slide_content': [],
                'tables': [],
                'formulas': [],
                'errors': []
            }
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                try:
                    slide_results = self._process_slide_advanced(
                        slide, doc_id, slide_idx, push_record_func
                    )
                    
                    # Aggregate results
                    processing_results['images'].extend(slide_results.get('images', []))
                    processing_results['slide_content'].extend(slide_results.get('slide_content', []))
                    processing_results['tables'].extend(slide_results.get('tables', []))
                    processing_results['formulas'].extend(slide_results.get('formulas', []))
                    processing_results['processed_slides'] += 1
                    
                except Exception as e:
                    error_msg = f"Error processing slide {slide_idx}: {str(e)}"
                    self.logger.error(error_msg)
                    processing_results['errors'].append(error_msg)
            
            return processing_results
            
        except Exception as e:
            error_msg = f"Error processing PPTX {path}: {str(e)}"
            self.logger.error(error_msg)
            return {'errors': [error_msg]}
    
    def _process_slide_advanced(self, slide, doc_id: str, slide_idx: int, push_record_func) -> Dict[str, Any]:
        """Process a single slide with comprehensive analysis"""
        
        slide_results = {
            'images': [],
            'slide_content': [],
            'tables': [],
            'formulas': []
        }
        
        # Create output directory
        out_dir = OUT_DIR / doc_id / f"slide{slide_idx}"
        out_dir.mkdir(parents=True, exist_ok=True)
        
        # Extract comprehensive slide content
        slide_content = self._extract_comprehensive_slide_content(slide)
        
        # Generate slide screenshot for full context OCR
        slide_image = self._generate_slide_image(slide, out_dir, slide_idx)
        
        # Process slide image with OCR for additional text extraction
        ocr_results = None
        if slide_image:
            try:
                slide_array = np.array(Image.open(slide_image))
                ocr_results = self.ocr.process_document_page(slide_array)
            except Exception as e:
                self.logger.warning(f"OCR processing failed for slide {slide_idx}: {e}")
        
        # Process embedded images
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_results = self._process_embedded_image(
                        shape, slide, slide_content, ocr_results, 
                        out_dir, doc_id, slide_idx, img_idx, push_record_func
                    )
                    slide_results['images'].append(image_results)
                    img_idx += 1
                except Exception as e:
                    self.logger.warning(f"Failed to process image {img_idx} on slide {slide_idx}: {e}")
        
        # Process tables detected by OCR
        if ocr_results:
            for table_idx, table_result in enumerate(ocr_results.get('table_regions', [])):
                try:
                    table_results = self._process_ocr_table(
                        table_result, slide_array, slide_content, 
                        out_dir, doc_id, slide_idx, table_idx, push_record_func
                    )
                    slide_results['tables'].append(table_results)
                except Exception as e:
                    self.logger.warning(f"Failed to process table {table_idx} on slide {slide_idx}: {e}")
        
        # Process formulas detected by OCR
        if ocr_results:
            for formula_idx, formula_result in enumerate(ocr_results.get('formula_regions', [])):
                try:
                    formula_results = self._process_ocr_formula(
                        formula_result, slide_array, slide_content, 
                        out_dir, doc_id, slide_idx, formula_idx, push_record_func
                    )
                    slide_results['formulas'].append(formula_results)
                except Exception as e:
                    self.logger.warning(f"Failed to process formula {formula_idx} on slide {slide_idx}: {e}")
        
        # Store comprehensive slide content
        slide_content_record = {
            'slide_number': slide_idx,
            'content': slide_content,
            'ocr_extracted_text': self._extract_ocr_text(ocr_results) if ocr_results else "",
            'image_count': img_idx,
            'table_count': len(slide_results['tables']),
            'formula_count': len(slide_results['formulas'])
        }
        
        # Push slide content as a record
        push_record_func(
            doc_id,
            f"slide{slide_idx}_content",
            f"Slide {slide_idx} content: {slide_content['title'] or 'Untitled'} - {slide_content['summary'][:100]}",
            slide_image,
            raw_text=json.dumps(slide_content, indent=2)
        )
        
        slide_results['slide_content'].append(slide_content_record)
        
        return slide_results
    
    def _extract_comprehensive_slide_content(self, slide) -> Dict[str, Any]:
        """Extract all text content from a slide"""
        content = {
            'title': '',
            'content_text': [],
            'notes': '',
            'alt_texts': [],
            'summary': '',
            'text_shapes': [],
            'layout_info': {}
        }
        
        # Extract title
        if slide.shapes.title:
            content['title'] = slide.shapes.title.text
        
        # Extract all text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape_info = {
                    'text': shape.text,
                    'shape_type': str(shape.shape_type),
                    'position': {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    }
                }
                content['text_shapes'].append(shape_info)
                content['content_text'].append(shape.text)
            
            # Extract alt text from images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, 'alternative_text'):
                if shape.alternative_text:
                    content['alt_texts'].append(shape.alternative_text)
        
        # Extract notes
        if slide.has_notes_slide:
            try:
                content['notes'] = slide.notes_slide.notes_text_frame.text
            except:
                content['notes'] = ''
        
        # Create summary
        all_text = []
        if content['title']:
            all_text.append(content['title'])
        all_text.extend(content['content_text'])
        if content['notes']:
            all_text.append(content['notes'])
        
        content['summary'] = ' '.join(all_text)[:500]  # First 500 chars as summary
        
        # Layout information
        content['layout_info'] = {
            'slide_layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
            'shape_count': len(slide.shapes),
            'text_shape_count': len(content['text_shapes']),
            'has_title': bool(content['title']),
            'has_notes': bool(content['notes'])
        }
        
        return content
    
    def _generate_slide_image(self, slide, out_dir: Path, slide_idx: int) -> Optional[Path]:
        """Generate an image of the entire slide for OCR processing"""
        try:
            # Note: This is a simplified approach. In practice, you might need
            # python-pptx-interface or other libraries for slide rendering
            # For now, we'll create a placeholder that could be enhanced
            
            slide_image_path = out_dir / f"slide{slide_idx}_full.png"
            
            # Create a basic slide representation
            # This would ideally render the actual slide content
            img = Image.new('RGB', (1024, 768), 'white')
            draw = ImageDraw.Draw(img)
            
            # Add slide content as text overlay (simplified)
            y_offset = 50
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Wrap text and draw
                    lines = self._wrap_text(shape.text, 80)
                    for line in lines[:10]:  # Limit to 10 lines
                        draw.text((50, y_offset), line, fill='black')
                        y_offset += 25
                    y_offset += 20
            
            img.save(slide_image_path, 'PNG')
            return slide_image_path
            
        except Exception as e:
            self.logger.warning(f"Failed to generate slide image: {e}")
            return None
    
    def _wrap_text(self, text: str, width: int) -> List[str]:
        """Simple text wrapping"""
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 <= width:
                current_line.append(word)
                current_length += len(word) + 1
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
                current_length = len(word)
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines
    
    def _process_embedded_image(self, shape, slide, slide_content: Dict, ocr_results, 
                               out_dir: Path, doc_id: str, slide_idx: int, img_idx: int, 
                               push_record_func) -> Dict[str, Any]:
        """Process an embedded image with comprehensive context"""
        
        # Extract and save image
        img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
        img_path = out_dir / f"img{img_idx}.png"
        img.save(img_path)
        
        # Get comprehensive caption using multiple methods
        caption = self._get_comprehensive_caption(
            shape, slide_content, img_path, img_idx
        )
        
        # Create comprehensive context
        context_info = {
            'slide_number': slide_idx,
            'slide_title': slide_content.get('title', ''),
            'slide_summary': slide_content.get('summary', ''),
            'slide_notes': slide_content.get('notes', ''),
            'image_alt_text': shape.alternative_text if hasattr(shape, 'alternative_text') else '',
            'image_position': {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            },
            'surrounding_text': self._get_surrounding_text(shape, slide),
            'slide_context': slide_content.get('content_text', [])
        }
        
        # Perform OCR on the image itself
        img_ocr_results = self._ocr_analyze_image(img_path)
        if img_ocr_results:
            context_info['image_ocr_text'] = img_ocr_results.get('text', '')
            context_info['image_contains_text'] = bool(img_ocr_results.get('text', '').strip())
        
        # Enhanced description combining all context
        enhanced_description = self._create_enhanced_description(caption, context_info)
        
        # Push to record system
        push_record_func(
            doc_id,
            f"s{slide_idx}_img{img_idx}",
            enhanced_description,
            img_path,
            raw_text=json.dumps(context_info, indent=2)
        )
        
        return {
            'path': str(img_path),
            'caption': enhanced_description,
            'context': context_info,
            'slide': slide_idx,
            'index': img_idx
        }
    
    def _get_comprehensive_caption(self, shape, slide_content: Dict, img_path: Path, img_idx: int) -> str:
        """Generate comprehensive caption using multiple methods"""
        
        captions = []
        
        # Method 1: Alternative text
        if hasattr(shape, 'alternative_text') and shape.alternative_text:
            captions.append(('alt_text', shape.alternative_text, 0.9))
        
        # Method 2: AI-generated caption
        try:
            ai_caption = auto_caption(str(img_path))
            if ai_caption and len(ai_caption.strip()) > 5:
                captions.append(('ai_caption', ai_caption, 0.8))
        except Exception as e:
            self.logger.warning(f"AI captioning failed: {e}")
        
        # Method 3: OCR on image
        try:
            img_array = np.array(Image.open(img_path))
            ocr_result = self.ocr.extract_text_multilingual(img_array)
            if ocr_result.confidence > 0.7 and len(ocr_result.text.strip()) > 5:
                captions.append(('ocr_text', f"Image containing text: {ocr_result.text[:100]}", 0.7))
        except Exception as e:
            self.logger.warning(f"Image OCR failed: {e}")
        
        # Method 4: Context-based caption
        slide_title = slide_content.get('title', '')
        if slide_title:
            context_caption = f"Image from slide '{slide_title}'"
            captions.append(('context', context_caption, 0.6))
        
        # Select best caption or combine
        if captions:
            # Sort by confidence and return the best
            captions.sort(key=lambda x: x[2], reverse=True)
            return captions[0][1]
        else:
            return f"Image {img_idx + 1} from slide {slide_content.get('title', 'Untitled')}"
    
    def _get_surrounding_text(self, target_shape, slide) -> str:
        """Get text from shapes near the target image"""
        surrounding_text = []
        
        target_left = target_shape.left
        target_top = target_shape.top
        target_right = target_left + target_shape.width
        target_bottom = target_top + target_shape.height
        
        # Find nearby text shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape_left = shape.left
                shape_top = shape.top
                shape_right = shape_left + shape.width
                shape_bottom = shape_top + shape.height
                
                # Check if shapes are nearby (simple proximity check)
                horizontal_overlap = not (target_right < shape_left or shape_right < target_left)
                vertical_overlap = not (target_bottom < shape_top or shape_bottom < target_top)
                
                # Consider text that's close to the image
                nearby_threshold = 100000  # Adjust based on typical slide dimensions
                horizontal_distance = min(abs(target_left - shape_right), abs(shape_left - target_right))
                vertical_distance = min(abs(target_top - shape_bottom), abs(shape_top - target_bottom))
                
                if (horizontal_overlap or vertical_overlap or 
                    horizontal_distance < nearby_threshold or 
                    vertical_distance < nearby_threshold):
                    surrounding_text.append(shape.text)
        
        return ' '.join(surrounding_text)
    
    def _ocr_analyze_image(self, img_path: Path) -> Optional[Dict[str, Any]]:
        """Perform OCR analysis on a single image"""
        try:
            img_array = np.array(Image.open(img_path))
            results = self.ocr.process_document_page(img_array)
            
            # Extract text from all regions
            all_text = []
            for region in results.get('text_regions', []):
                if region.confidence > 0.5:
                    all_text.append(region.text)
            
            return {
                'text': ' '.join(all_text),
                'has_tables': len(results.get('table_regions', [])) > 0,
                'has_formulas': len(results.get('formula_regions', [])) > 0,
                'confidence': max([r.confidence for r in results.get('text_regions', [])], default=0.0)
            }
        except Exception as e:
            self.logger.warning(f"OCR analysis failed for {img_path}: {e}")
            return None
    
    def _create_enhanced_description(self, base_caption: str, context_info: Dict) -> str:
        """Create an enhanced description combining caption and context"""
        
        description_parts = [base_caption]
        
        # Add slide context
        if context_info.get('slide_title'):
            description_parts.append(f"From slide: '{context_info['slide_title']}'")
        
        # Add image text if present
        if context_info.get('image_contains_text') and context_info.get('image_ocr_text'):
            text_preview = context_info['image_ocr_text'][:100]
            description_parts.append(f"Contains text: {text_preview}")
        
        # Add surrounding context
        if context_info.get('surrounding_text'):
            context_preview = context_info['surrounding_text'][:100]
            description_parts.append(f"Context: {context_preview}")
        
        # Add position information
        pos = context_info.get('image_position', {})
        if pos:
            description_parts.append(f"Position: slide coordinates ({pos.get('left', 0)}, {pos.get('top', 0)})")
        
        return ' | '.join(description_parts)
    
    def _process_ocr_table(self, table_result: Dict, slide_image: np.ndarray, slide_content: Dict,
                          out_dir: Path, doc_id: str, slide_idx: int, table_idx: int, 
                          push_record_func) -> Dict[str, Any]:
        """Process a table detected by OCR"""
        
        # Save table region as image
        bbox = table_result.get('bbox', (0, 0, 100, 100))
        x1, y1, x2, y2 = bbox
        table_img = slide_image[y1:y2, x1:x2]
        
        table_path = out_dir / f"table{table_idx}.png"
        Image.fromarray(table_img).save(table_path, "PNG")
        
        # Create table description
        table_data = table_result.get('data', [])
        description = f"Table from slide {slide_idx}: {self._describe_table_data(table_data)}"
        
        # Add slide context
        if slide_content.get('title'):
            description += f" | Slide: {slide_content['title']}"
        
        push_record_func(
            doc_id,
            f"s{slide_idx}_table{table_idx}",
            description,
            table_path,
            raw_text=json.dumps(table_data)
        )
        
        return {
            'path': str(table_path),
            'description': description,
            'data': table_data,
            'slide': slide_idx,
            'index': table_idx
        }
    
    def _process_ocr_formula(self, formula_result, slide_image: np.ndarray, slide_content: Dict,
                            out_dir: Path, doc_id: str, slide_idx: int, formula_idx: int,
                            push_record_func) -> Dict[str, Any]:
        """Process a mathematical formula detected by OCR"""
        
        # Save formula region as image
        bbox = formula_result.bbox
        x1, y1, x2, y2 = bbox
        formula_img = slide_image[max(0, y1-10):y2+10, max(0, x1-10):x2+10]  # Add padding
        
        formula_path = out_dir / f"formula{formula_idx}.png"
        Image.fromarray(formula_img).save(formula_path, "PNG")
        
        # Create formula description
        latex = formula_result.metadata.get('latex', formula_result.text)
        description = f"Mathematical formula from slide {slide_idx}: {latex}"
        
        # Add slide context
        if slide_content.get('title'):
            description += f" | Slide: {slide_content['title']}"
        
        push_record_func(
            doc_id,
            f"s{slide_idx}_formula{formula_idx}",
            description,
            formula_path,
            raw_text=formula_result.text
        )
        
        return {
            'path': str(formula_path),
            'description': description,
            'text': formula_result.text,
            'latex': latex,
            'slide': slide_idx,
            'index': formula_idx
        }
    
    def _describe_table_data(self, table_data: List[List[str]]) -> str:
        """Create a description for table data"""
        if not table_data:
            return "Empty table"
        
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        description = f"{rows} rows × {cols} columns"
        
        # Add header info if available
        if rows > 0 and cols > 0:
            headers = [cell for cell in table_data[0] if cell.strip()]
            if headers:
                description += f", headers: {', '.join(headers[:3])}"
                if len(headers) > 3:
                    description += "..."
        
        return description
    
    def _extract_ocr_text(self, ocr_results: Dict) -> str:
        """Extract all text from OCR results"""
        all_text = []
        
        for region in ocr_results.get('text_regions', []):
            if region.confidence > 0.5:
                all_text.append(region.text)
        
        for region in ocr_results.get('handwriting_regions', []):
            if region.confidence > 0.3:
                all_text.append(region.text)
        
        return ' '.join(all_text)


# Backward compatibility function
def process_pptx(path: Path, push_record):
    """
    Enhanced PPTX processing function with Revolutionary OCR
    Maintains compatibility with existing interface while adding advanced features
    """
    try:
        # Initialize enhanced processor
        config = get_revolutionary_ocr_config()
        processor = EnhancedPPTXProcessor(config['ocr'])
        
        # Process with advanced capabilities
        results = processor.process_pptx_advanced(path, push_record)
        
        # Log processing summary
        logger.info(f"Processed PPTX {path.name}: "
                   f"{results.get('processed_slides', 0)}/{results.get('total_slides', 0)} slides, "
                   f"{len(results.get('images', []))} images, "
                   f"{len(results.get('tables', []))} tables, "
                   f"{len(results.get('formulas', []))} formulas")
        
        if results.get('errors'):
            logger.warning(f"Errors during processing: {results['errors']}")
        
        return results
        
    except Exception as e:
        logger.error(f"Failed to process PPTX {path}: {e}")
        # Fallback to basic processing
        return _process_pptx_basic(path, push_record)


def _process_pptx_basic(path: Path, push_record):
    """Fallback to basic PPTX processing if advanced processing fails"""
    try:
        prs = Presentation(path)
        doc_id = path.stem
        
        for s_idx, slide in enumerate(prs.slides, 1):
            title = slide.shapes.title.text if slide.shapes.title else ""
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            out_dir = OUT_DIR / doc_id / f"slide{s_idx}"
            out_dir.mkdir(parents=True, exist_ok=True)
            
            img_i = 0
            for sh in slide.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                    p = out_dir / f"img{img_i}.png"
                    img.save(p)
                    caption = sh.alternative_text or title or notes or auto_caption(str(p))
                    push_record(doc_id, f"s{s_idx}_img{img_i}", caption, p, raw_text=title+"\n"+notes)
                    img_i += 1
                    
    except Exception as e:
        logger.error(f"Basic PPTX processing also failed: {e}")
        raise
