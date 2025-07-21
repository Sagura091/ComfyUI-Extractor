"""
Comprehensive Test Suite for Enhanced Document Processing
Tests PPTX, XLSX, and Markdown processing with Revolutionary OCR
"""

import sys
import os
import tempfile
import json
import time
from pathlib import Path
from typing import List, Dict, Any

# Add the extractor directory to the Python path
sys.path.insert(0, str(Path(__file__).parent))

try:
    from enhanced_pptx_utils import EnhancedPPTXProcessor
    from enhanced_xlsx_utils import EnhancedXLSXProcessor
    from enhanced_md_utils import EnhancedMarkdownProcessor
    from ocr_config import get_revolutionary_ocr_config, validate_config
    from revolutionary_ocr import create_revolutionary_ocr
    print("✅ All enhanced utilities imported successfully")
except ImportError as e:
    print(f"❌ Import error: {e}")
    sys.exit(1)

class TestResults:
    """Track test results and metrics"""
    
    def __init__(self):
        self.tests_run = 0
        self.tests_passed = 0
        self.tests_failed = 0
        self.errors = []
        self.warnings = []
        self.performance_metrics = {}
    
    def add_test(self, name: str, passed: bool, duration: float = 0, error: str = None):
        self.tests_run += 1
        if passed:
            self.tests_passed += 1
            print(f"✅ {name} ({duration:.2f}s)")
        else:
            self.tests_failed += 1
            self.errors.append(f"{name}: {error}")
            print(f"❌ {name}: {error}")
        
        self.performance_metrics[name] = duration
    
    def add_warning(self, message: str):
        self.warnings.append(message)
        print(f"⚠️  {message}")
    
    def print_summary(self):
        print("\n" + "="*80)
        print("TEST SUMMARY")
        print("="*80)
        print(f"Tests Run: {self.tests_run}")
        print(f"Passed: {self.tests_passed}")
        print(f"Failed: {self.tests_failed}")
        print(f"Success Rate: {(self.tests_passed/self.tests_run)*100:.1f}%")
        
        if self.warnings:
            print(f"\nWarnings: {len(self.warnings)}")
            for warning in self.warnings:
                print(f"  - {warning}")
        
        if self.errors:
            print(f"\nErrors: {len(self.errors)}")
            for error in self.errors:
                print(f"  - {error}")
        
        print(f"\nPerformance Metrics:")
        for test_name, duration in self.performance_metrics.items():
            print(f"  {test_name}: {duration:.2f}s")

def test_configuration():
    """Test OCR configuration validation"""
    results = TestResults()
    
    # Test 1: Configuration validation
    start_time = time.time()
    try:
        config_valid = validate_config()
        results.add_test("Configuration Validation", config_valid, time.time() - start_time)
    except Exception as e:
        results.add_test("Configuration Validation", False, time.time() - start_time, str(e))
    
    # Test 2: Configuration loading
    start_time = time.time()
    try:
        config = get_revolutionary_ocr_config()
        has_required_sections = all(section in config for section in 
                                  ['ocr', 'language_detection', 'table_recognition', 'formula_recognition'])
        results.add_test("Configuration Loading", has_required_sections, time.time() - start_time)
        
        if has_required_sections:
            print(f"   Languages: {len(config['ocr']['languages'])}")
            print(f"   Features: Tables={config['ocr']['enable_table_recognition']}, "
                  f"Formulas={config['ocr']['enable_formula_recognition']}, "
                  f"Handwriting={config['ocr']['enable_handwriting']}")
    except Exception as e:
        results.add_test("Configuration Loading", False, time.time() - start_time, str(e))
    
    return results

def test_ocr_initialization():
    """Test Revolutionary OCR system initialization"""
    results = TestResults()
    
    # Test OCR creation
    start_time = time.time()
    try:
        config = get_revolutionary_ocr_config()
        ocr = create_revolutionary_ocr(config['ocr'])
        results.add_test("OCR Initialization", True, time.time() - start_time)
        
        # Test basic functionality
        import numpy as np
        test_image = np.ones((100, 100, 3), dtype=np.uint8) * 255  # White image
        
        start_time = time.time()
        result = ocr.extract_text_multilingual(test_image)
        results.add_test("OCR Basic Processing", True, time.time() - start_time)
        
    except Exception as e:
        results.add_test("OCR Initialization", False, time.time() - start_time, str(e))
    
    return results

def create_test_pptx():
    """Create a simple test PPTX file"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import io
        
        prs = Presentation()
        
        # Slide 1: Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        slide.shapes.placeholders[1].text = "Enhanced Processing Test"
        
        # Slide 2: Content with image
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide with Image"
        
        # Create a simple test image
        img = Image.new('RGB', (200, 100), color='lightblue')
        img_io = io.BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)
        
        # Add image to slide
        left = Inches(1)
        top = Inches(2)
        slide.shapes.add_picture(img_io, left, top, width=Inches(3))
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        prs.save(temp_file.name)
        return Path(temp_file.name)
        
    except Exception as e:
        print(f"Failed to create test PPTX: {e}")
        return None

def create_test_xlsx():
    """Create a simple test XLSX file"""
    try:
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image
        from PIL import Image as PILImage
        import io
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Data"
        
        # Add sample data
        ws['A1'] = "Month"
        ws['B1'] = "Revenue"
        ws['C1'] = "Profit"
        
        data = [
            ["January", 1000, 200],
            ["February", 1200, 250],
            ["March", 1100, 220]
        ]
        
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        
        # Create and add a simple image
        img = PILImage.new('RGB', (150, 100), color='lightgreen')
        img_io = io.BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
        wb.save(temp_file.name)
        return Path(temp_file.name)
        
    except Exception as e:
        print(f"Failed to create test XLSX: {e}")
        return None

def create_test_markdown():
    """Create a simple test Markdown file"""
    try:
        content = """# Test Document

This is a test markdown document for enhanced processing.

## Section 1: Introduction

This section contains some introductory text and an image.

![Test Image](https://via.placeholder.com/300x200/blue/white?text=Test+Image)

## Section 2: Data Table

| Column 1 | Column 2 | Column 3 |
|----------|----------|----------|
| Value 1  | Value 2  | Value 3  |
| Data A   | Data B   | Data C   |

## Section 3: Code Example

```python
def hello_world():
    print("Hello, World!")
    return "Enhanced processing test"
```

## Section 4: Mathematical Formula

The quadratic formula is: $x = \\frac{-b \\pm \\sqrt{b^2-4ac}}{2a}$

End of document.
"""
        
        temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8')
        temp_file.write(content)
        temp_file.close()
        return Path(temp_file.name)
        
    except Exception as e:
        print(f"Failed to create test Markdown: {e}")
        return None

def test_pptx_processing():
    """Test enhanced PPTX processing"""
    results = TestResults()
    
    # Create test PPTX
    test_file = create_test_pptx()
    if not test_file:
        results.add_test("PPTX Test File Creation", False, 0, "Could not create test file")
        return results
    
    try:
        # Test processor initialization
        start_time = time.time()
        config = get_revolutionary_ocr_config()
        processor = EnhancedPPTXProcessor(config['ocr'])
        results.add_test("PPTX Processor Initialization", True, time.time() - start_time)
        
        # Test processing
        start_time = time.time()
        processed_records = []
        
        def mock_push_record(doc_id, part_id, caption, img_path, raw_text=""):
            processed_records.append({
                'doc_id': doc_id,
                'part_id': part_id,
                'caption': caption,
                'img_path': str(img_path) if img_path else None,
                'raw_text': raw_text[:100]
            })
        
        processing_results = processor.process_pptx_advanced(test_file, mock_push_record)
        duration = time.time() - start_time
        
        # Validate results
        success = (
            processing_results.get('processed_slides', 0) > 0 and
            'errors' in processing_results and
            len(processed_records) > 0
        )
        
        results.add_test("PPTX Processing", success, duration)
        
        if success:
            print(f"   Processed {processing_results.get('processed_slides', 0)} slides")
            print(f"   Generated {len(processed_records)} records")
            print(f"   Found {len(processing_results.get('images', []))} images")
        
    except Exception as e:
        results.add_test("PPTX Processing", False, 0, str(e))
    finally:
        # Cleanup
        if test_file.exists():
            test_file.unlink()
    
    return results

def test_xlsx_processing():
    """Test enhanced XLSX processing"""
    results = TestResults()
    
    # Create test XLSX
    test_file = create_test_xlsx()
    if not test_file:
        results.add_test("XLSX Test File Creation", False, 0, "Could not create test file")
        return results
    
    try:
        # Test processor initialization
        start_time = time.time()
        config = get_revolutionary_ocr_config()
        processor = EnhancedXLSXProcessor(config['ocr'])
        results.add_test("XLSX Processor Initialization", True, time.time() - start_time)
        
        # Test processing
        start_time = time.time()
        processed_records = []
        
        def mock_push_record(doc_id, part_id, caption, img_path, raw_text=""):
            processed_records.append({
                'doc_id': doc_id,
                'part_id': part_id,
                'caption': caption,
                'img_path': str(img_path) if img_path else None,
                'raw_text': raw_text[:100]
            })
        
        processing_results = processor.process_xlsx_advanced(test_file, mock_push_record)
        duration = time.time() - start_time
        
        # Validate results
        success = (
            processing_results.get('processed_worksheets', 0) > 0 and
            'errors' in processing_results and
            len(processed_records) > 0
        )
        
        results.add_test("XLSX Processing", success, duration)
        
        if success:
            print(f"   Processed {processing_results.get('processed_worksheets', 0)} worksheets")
            print(f"   Generated {len(processed_records)} records")
            print(f"   Found {len(processing_results.get('images', []))} images")
        
    except Exception as e:
        results.add_test("XLSX Processing", False, 0, str(e))
    finally:
        # Cleanup
        if test_file.exists():
            test_file.unlink()
    
    return results

def test_markdown_processing():
    """Test enhanced Markdown processing"""
    results = TestResults()
    
    # Create test Markdown
    test_file = create_test_markdown()
    if not test_file:
        results.add_test("Markdown Test File Creation", False, 0, "Could not create test file")
        return results
    
    try:
        # Test processor initialization
        start_time = time.time()
        config = get_revolutionary_ocr_config()
        processor = EnhancedMarkdownProcessor(config['ocr'])
        results.add_test("Markdown Processor Initialization", True, time.time() - start_time)
        
        # Test processing
        start_time = time.time()
        processed_records = []
        
        def mock_push_record(doc_id, part_id, caption, img_path, raw_text=""):
            processed_records.append({
                'doc_id': doc_id,
                'part_id': part_id,
                'caption': caption,
                'img_path': str(img_path) if img_path else None,
                'raw_text': raw_text[:100]
            })
        
        processing_results = processor.process_markdown_advanced(test_file, mock_push_record)
        duration = time.time() - start_time
        
        # Validate results
        success = (
            'document_structure' in processing_results and
            'errors' in processing_results and
            len(processed_records) > 0
        )
        
        results.add_test("Markdown Processing", success, duration)
        
        if success:
            print(f"   Generated {len(processed_records)} records")
            print(f"   Found {len(processing_results.get('images', []))} images")
            print(f"   Found {len(processing_results.get('code_blocks', []))} code blocks")
            print(f"   Found {len(processing_results.get('tables', []))} tables")
        
    except Exception as e:
        results.add_test("Markdown Processing", False, 0, str(e))
    finally:
        # Cleanup
        if test_file.exists():
            test_file.unlink()
    
    return results

def run_all_tests():
    """Run comprehensive test suite"""
    print("🚀 Starting Enhanced Document Processing Test Suite")
    print("="*80)
    
    all_results = TestResults()
    
    # Test 1: Configuration
    print("\n📋 Testing Configuration...")
    config_results = test_configuration()
    all_results.tests_run += config_results.tests_run
    all_results.tests_passed += config_results.tests_passed
    all_results.tests_failed += config_results.tests_failed
    all_results.errors.extend(config_results.errors)
    all_results.warnings.extend(config_results.warnings)
    all_results.performance_metrics.update(config_results.performance_metrics)
    
    # Test 2: OCR Initialization
    print("\n🔍 Testing OCR Initialization...")
    ocr_results = test_ocr_initialization()
    all_results.tests_run += ocr_results.tests_run
    all_results.tests_passed += ocr_results.tests_passed
    all_results.tests_failed += ocr_results.tests_failed
    all_results.errors.extend(ocr_results.errors)
    all_results.warnings.extend(ocr_results.warnings)
    all_results.performance_metrics.update(ocr_results.performance_metrics)
    
    # Test 3: PPTX Processing
    print("\n📊 Testing PPTX Processing...")
    pptx_results = test_pptx_processing()
    all_results.tests_run += pptx_results.tests_run
    all_results.tests_passed += pptx_results.tests_passed
    all_results.tests_failed += pptx_results.tests_failed
    all_results.errors.extend(pptx_results.errors)
    all_results.warnings.extend(pptx_results.warnings)
    all_results.performance_metrics.update(pptx_results.performance_metrics)
    
    # Test 4: XLSX Processing
    print("\n📈 Testing XLSX Processing...")
    xlsx_results = test_xlsx_processing()
    all_results.tests_run += xlsx_results.tests_run
    all_results.tests_passed += xlsx_results.tests_passed
    all_results.tests_failed += xlsx_results.tests_failed
    all_results.errors.extend(xlsx_results.errors)
    all_results.warnings.extend(xlsx_results.warnings)
    all_results.performance_metrics.update(xlsx_results.performance_metrics)
    
    # Test 5: Markdown Processing
    print("\n📝 Testing Markdown Processing...")
    md_results = test_markdown_processing()
    all_results.tests_run += md_results.tests_run
    all_results.tests_passed += md_results.tests_passed
    all_results.tests_failed += md_results.tests_failed
    all_results.errors.extend(md_results.errors)
    all_results.warnings.extend(md_results.warnings)
    all_results.performance_metrics.update(md_results.performance_metrics)
    
    # Print final summary
    all_results.print_summary()
    
    # Recommendations
    print("\n💡 RECOMMENDATIONS:")
    if all_results.tests_passed == all_results.tests_run:
        print("🎉 All tests passed! The enhanced processing system is ready for production.")
    elif all_results.tests_passed >= all_results.tests_run * 0.8:
        print("✅ Most tests passed. Check warnings and errors before production deployment.")
    else:
        print("⚠️  Several tests failed. Review configuration and dependencies before deployment.")
    
    print("\n📚 Next Steps:")
    print("1. Install missing dependencies: pip install -r requirements.txt")
    print("2. Review configuration in ocr_config.py")
    print("3. Test with real documents using the API")
    print("4. Monitor performance and adjust settings as needed")
    
    return all_results.tests_passed == all_results.tests_run

if __name__ == "__main__":
    success = run_all_tests()
    sys.exit(0 if success else 1)
