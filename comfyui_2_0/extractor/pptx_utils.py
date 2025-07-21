from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from PIL import Image
import io
from .caption_api import auto_caption
from .settings import OUT_DIR

def process_pptx(path: Path, push_record):
    prs = Presentation(path); doc_id = path.stem
    for s_idx, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text if slide.shapes.title else ""
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        out_dir = OUT_DIR / doc_id / f"slide{s_idx}"; out_dir.mkdir(parents=True, exist_ok=True)
        img_i = 0
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                p = out_dir / f"img{img_i}.png"; img.save(p)
                caption = sh.alternative_text or title or notes or auto_caption(p)
                push_record(doc_id, f"s{s_idx}_img{img_i}", caption, p, raw_text=title+"\n"+notes)
                img_i += 1
